"""Generate the PowerPoint and Word templates from the HBRLRG design system.

The HTML templates in this repository are the source of truth — they are text,
so a change to them shows up in a diff. Office files are binary, so they are
built from the same tokens at release time and attached to the release rather
than committed. See DEPLOY.md §6.

    python3 tools/build_office.py [outdir]

Needs python-pptx, python-docx, Pillow. Fonts are embedded into the .pptx when
a directory of Archivo TTFs is passed as HBRLRG_FONT_DIR.
"""
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import docx
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Cm, Mm, Pt as DocPt, RGBColor as DocRGB

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"

# --- design tokens ---------------------------------------------------------
NAVY = RGBColor(0x14, 0x24, 0x3F)
TEAL = RGBColor(0x0E, 0xA7, 0x9A)
INK = RGBColor(0x2B, 0x2F, 0x36)
MUTED = RGBColor(0x6B, 0x72, 0x80)
SURFACE = RGBColor(0xF6, 0xF7, 0xF9)
HAIRLINE = RGBColor(0xE4, 0xE7, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ON_NAVY = RGBColor(0x93, 0xA2, 0xBC)
ON_NAVY_DIM = RGBColor(0x6E, 0x82, 0xA3)
NAVY_RULE = RGBColor(0x2C, 0x40, 0x66)
FOOT = RGBColor(0x9A, 0xA1, 0xAB)

LATIN = "Archivo"
HANGUL = "Pretendard"

# The master is 1920x1080 on a 13.333in stage: 144px per inch, and a pixel of
# type is half a point. Every number below is the px value from the HTML.
PXI = 1.0 / 144.0
W_PX, H_PX = 1920, 1080
MARGIN = 120
BODY_W = W_PX - 2 * MARGIN

CONTACT = "HBRL Research Group  ·  hbrlee@unist.ac.kr  ·  hbrl-research.group"
NAME = "Han-Bo-Ram Lee"
AFFIL = ("Professor, Graduate School of Semiconductor Materials "
         "& Devices Engineering, UNIST")
EDITOR = "Executive Editor, Chemistry of Materials\nACS Publications"


def px(v):
    return Inches(v * PXI)


def ppt(v):
    return Pt(v / 2.0)


# --- pptx helpers ----------------------------------------------------------
def rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, body, size, colour, *, bold=False, spacing=None,
         align=PP_ALIGN.LEFT, tracking=None, wrap=True, anchor=None):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = LATIN
        r.font.size = ppt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        # East-Asian font, so Korean typed into the template lands on Pretendard
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {"typeface": HANGUL})
        rPr.append(ea)
        if tracking is not None:
            rPr.set("spc", str(int(tracking * size / 2 * 100)))
    return box


def picture(slide, name, x, y, h):
    from PIL import Image
    path = ASSETS / name
    with Image.open(path) as im:
        w = h * im.width / im.height
    return slide.shapes.add_picture(str(path), px(x), px(y), px(w), px(h))


def header(slide, eyebrow, heading, subtitle=None):
    """Eyebrow, heading, and the teal-capped divider — the lab's signature."""
    y = 56
    if eyebrow:
        text(slide, MARGIN, y, BODY_W, 30, eyebrow, 24, TEAL,
             bold=False, tracking=0.14)
        y += 48
    text(slide, MARGIN, y, BODY_W, 90, heading, 60, NAVY, bold=True, spacing=1.1)
    y += 90
    rect(slide, MARGIN, y + 1.5, BODY_W, 2, HAIRLINE)
    rect(slide, MARGIN, y, 88, 5, TEAL)
    if subtitle:
        text(slide, MARGIN, y + 26, BODY_W, 44, subtitle, 28, MUTED)
    return 268


def footer(slide, on_navy=False):
    picture(slide, "logo/hbrlrg-horizontal-reverse.png" if on_navy
            else "logo/hbrlrg-horizontal.png", MARGIN, 1004, 42)
    text(slide, W_PX - MARGIN - 900, 1010, 900, 30, CONTACT, 24,
         ON_NAVY_DIM if on_navy else FOOT, align=PP_ALIGN.RIGHT, wrap=False)


def bullets(slide, y, items):
    for line in items:
        rect(slide, MARGIN, y + 17, 14, 14, TEAL)
        text(slide, MARGIN + 36, y, BODY_W - 36, 60, line, 34, INK, spacing=1.4)
        y += 72
    return y


def slot(slide, x, y, w, h, label):
    """A figure well — where a plot or micrograph gets pasted in."""
    box = rect(slide, x, y, w, h, SURFACE)
    box.line.color.rgb = HAIRLINE
    box.line.width = Pt(1)
    text(slide, x, y + h / 2 - 16, w, 32, label, 24, MUTED, align=PP_ALIGN.CENTER)
    return box


def build_pptx(out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = px(W_PX), px(H_PX)
    blank = prs.slide_layouts[6]

    def new(bg=WHITE):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, W_PX, H_PX, bg)
        return s

    # 01 Title -------------------------------------------------------------
    s = new(NAVY)
    picture(s, "logo/hbrlrg-horizontal-reverse.png", MARGIN, 96, 68)
    rect(s, MARGIN, 248, 88, 5, TEAL)
    text(s, MARGIN, 283, 1560, 240, "Presentation title goes here", 108, WHITE,
         bold=True, spacing=1.02)
    text(s, MARGIN, 560, 1500, 70, "One line that unpacks the title", 36, ON_NAVY)
    rect(s, MARGIN, 814, BODY_W, 1, NAVY_RULE)
    text(s, MARGIN, 834, 900, 50, NAME, 30, WHITE, bold=True)
    picture(s, "unist-emblem-onnavy.png", MARGIN, 890, 84)
    text(s, 228, 892, 560, 110, AFFIL, 24, ON_NAVY, spacing=1.35)
    picture(s, "cm-logo-onnavy.png", 880, 890, 84)
    text(s, 1060, 892, 640, 110, EDITOR, 24, ON_NAVY, spacing=1.35)
    text(s, MARGIN, 1008, 900, 40, "Venue  ·  Month 00, 2026", 24, ON_NAVY_DIM)
    s.notes_slide.notes_text_frame.text = (
        "표지. 제목은 26자 안팎에서 줄이 바뀝니다. 발표 장소와 날짜는 맨 아래 줄에."
    )

    # 02 Agenda ------------------------------------------------------------
    s = new()
    y = header(s, None, "Agenda")
    for i, item in enumerate(["Background", "ALD fundamentals",
                              "Process optimization", "Results", "Outlook"], 1):
        text(s, MARGIN, y, 80, 50, f"{i:02d}", 34, TEAL, bold=True)
        text(s, MARGIN + 90, y, BODY_W - 90, 50, item, 34, NAVY)
        y += 78
    footer(s)
    s.notes_slide.notes_text_frame.text = "목차. 번호는 본문 슬라이드 아이브로우와 짝을 맞춥니다."

    # 03 Section -----------------------------------------------------------
    s = new(NAVY)
    rect(s, MARGIN, 440, 88, 5, TEAL)
    text(s, MARGIN, 486, 1400, 60, "02", 32, TEAL, bold=True, tracking=0.14)
    text(s, MARGIN, 540, 1560, 160, "Section title", 104, WHITE, bold=True,
         spacing=1.05)
    s.notes_slide.notes_text_frame.text = "섹션 구분. 번호는 목차와 같게."

    # 04 Content -----------------------------------------------------------
    s = new()
    y = header(s, "02  ·  ALD FUNDAMENTALS", "Four steps of the self-limiting cycle")
    bullets(s, y, [
        "Precursor pulse saturates at a single monolayer",
        "First purge removes unreacted precursor and byproducts",
        "Reactant pulse completes the ligand exchange",
        "Second purge closes the cycle — thickness follows cycle count",
    ])
    footer(s)
    s.notes_slide.notes_text_frame.text = (
        "본문. 불릿은 상단 절반만 쓰고 아래는 그림 자리로 비워 둡니다."
    )

    # 05 Figure, full ------------------------------------------------------
    s = new()
    y = header(s, "03  ·  PROCESS OPTIMIZATION", "Wafer-scale uniformity")
    slot(s, MARGIN, y, BODY_W, 660, "Drop a plot, micrograph, or schematic")
    footer(s)

    # 06 Figure, two up ----------------------------------------------------
    s = new()
    y = header(s, "03  ·  PROCESS OPTIMIZATION", "Growth behavior by material")
    half = (BODY_W - 40) / 2
    slot(s, MARGIN, y, half, 600, "Left figure")
    slot(s, MARGIN + half + 40, y, half, 600, "Right figure")
    text(s, MARGIN, y + 620, half, 40, "Caption for the left figure", 22, MUTED)
    text(s, MARGIN + half + 40, y + 620, half, 40, "Caption for the right figure",
         22, MUTED)
    footer(s)

    # 07 Figure + text -----------------------------------------------------
    s = new()
    y = header(s, "03  ·  PROCESS OPTIMIZATION", "Uniformity across the wafer")
    slot(s, MARGIN, y, 980, 620, "Figure")
    bullets(s, y, [])
    tx = MARGIN + 1020
    for line in ["Thickness varies by less than 1% edge to edge",
                 "The window holds across three precursor systems",
                 "Cycle count, not exposure time, sets the thickness"]:
        text(s, tx, y, BODY_W - 1020, 120, line, 30, INK, spacing=1.4)
        y += 130
    footer(s)

    # 08 Data --------------------------------------------------------------
    s = new()
    y = header(s, "03  ·  PROCESS OPTIMIZATION", "Standard process conditions")
    cols = [520, 380, 380, 400]
    heads = ["Parameter", "HfO₂", "Al₂O₃", "Note"]
    rows = [["Deposition temperature", "250 °C", "200 °C", "ALD window"],
            ["Precursor", "TEMAH", "TMA", "—"],
            ["Cycles", "300", "250", "±0.8% reproducibility"],
            ["Substrate", "Si (100)", "Si (100)", "Room temperature load"]]
    x = MARGIN
    rect(s, MARGIN, y, BODY_W, 56, SURFACE)
    for w, h in zip(cols, heads):
        text(s, x + 16, y + 16, w - 32, 32, h, 24, MUTED, tracking=0.08)
        x += w
    y += 56
    for row in rows:
        x = MARGIN
        for w, cell in zip(cols, row):
            text(s, x + 16, y + 16, w - 32, 40, cell, 28, INK)
            x += w
        rect(s, MARGIN, y + 66, BODY_W, 1, HAIRLINE)
        y += 68
    text(s, MARGIN, y + 14, BODY_W, 40,
         "Table 1. Captions sit below the table, without terminal punctuation",
         22, MUTED)
    footer(s)

    # 09 Metrics -----------------------------------------------------------
    s = new(SURFACE)
    y = header(s, "03  ·  PROCESS OPTIMIZATION", "Three numbers that carry the talk")
    third = (BODY_W - 80) / 3
    for i, (fig, lab) in enumerate([("92%", "Deposition accuracy"),
                                    ("±0.8%", "Thickness reproducibility"),
                                    ("10+", "Cycles simulated end to end")]):
        x = MARGIN + i * (third + 40)
        rect(s, x, y, 88, 5, TEAL)
        text(s, x, y + 40, third, 130, fig, 108, NAVY, bold=True)
        text(s, x, y + 190, third, 60, lab, 26, MUTED, spacing=1.35)
    footer(s)

    # 10 Full-bleed image --------------------------------------------------
    s = new(NAVY)
    slot(s, 0, 0, W_PX, H_PX, "Drop an equipment or laboratory photograph")
    rect(s, MARGIN, 812, 88, 5, TEAL)
    text(s, MARGIN, 856, 1400, 100, "Caption over the image", 72, WHITE, bold=True)
    s.notes_slide.notes_text_frame.text = "전면 이미지. 사진을 넣으면 회색 자리를 덮습니다."

    # 11 Statement ---------------------------------------------------------
    s = new(TEAL)
    text(s, MARGIN, 400, BODY_W, 300,
         "One sentence the audience should leave with.", 84, WHITE, bold=True,
         spacing=1.15)
    s.notes_slide.notes_text_frame.text = "선언. 한 문장만. 강조색 위이므로 로고는 넣지 않습니다."

    # 12 Closing -----------------------------------------------------------
    s = new(NAVY)
    rect(s, MARGIN, 300, 88, 5, TEAL)
    text(s, MARGIN, 344, 1560, 130, "Thank you", 92, WHITE, bold=True)
    text(s, MARGIN, 486, 1400, 70, "Questions welcome", 36, ON_NAVY)
    rect(s, MARGIN, 700, BODY_W, 1, NAVY_RULE)
    text(s, MARGIN, 726, 900, 50, NAME, 30, WHITE, bold=True)
    picture(s, "unist-emblem-onnavy.png", MARGIN, 790, 84)
    text(s, 228, 792, 560, 110, AFFIL, 24, ON_NAVY, spacing=1.35)
    footer(s, on_navy=True)

    prs.save(out)
    return len(prs.slides.__iter__.__self__._sldIdLst)


# --- docx ------------------------------------------------------------------
D_NAVY, D_INK, D_MUTED = DocRGB(0x14, 0x24, 0x3F), DocRGB(0x2B, 0x2F, 0x36), DocRGB(0x6B, 0x72, 0x80)


def dset(run, size, colour, *, bold=False, latin=LATIN):
    run.font.name = latin
    run.font.size = DocPt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    rPr = run._r.get_or_add_rPr()
    rf = rPr.get_or_add_rFonts()
    rf.set(qn("w:eastAsia"), HANGUL)
    rf.set(qn("w:ascii"), latin)
    rf.set(qn("w:hAnsi"), latin)


def cell_border_bottom(cell, colour, size):
    """size is in eighths of a point."""
    tcPr = cell._tc.get_or_add_tcPr()
    borders = tcPr.makeelement(qn("w:tcBorders"), {})
    borders.append(borders.makeelement(qn("w:bottom"), {
        qn("w:val"): "single", qn("w:sz"): str(size),
        qn("w:space"): "0", qn("w:color"): colour,
    }))
    tcPr.append(borders)


def rule(doc):
    """The signature divider — a navy line with a teal cap at its left end.

    Drawn as a two-cell table because a paragraph border cannot carry two
    colours, and a glyph standing in for the cap is something a typist deletes
    by accident.
    """
    t = doc.add_table(rows=1, cols=2)
    t.autofit = False
    cap, line = t.rows[0].cells
    cap.width, line.width = Cm(1.2), Cm(16.2)
    cell_border_bottom(cap, "0EA79A", 16)
    cell_border_bottom(line, "14243F", 4)
    for c in (cap, line):
        c.paragraphs[0].paragraph_format.space_after = DocPt(0)
        c.paragraphs[0].runs and None
    doc.add_paragraph().paragraph_format.space_after = DocPt(4)
    return t


def doc_footer(doc, line, tracking_em):
    """Section footer, so it repeats on every page rather than sitting once
    at the end of the body."""
    para = doc.sections[0].footer.paragraphs[0]
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = para.add_run(line)
    dset(r, 7.5, DocRGB(0x9A, 0xA1, 0xAB))
    rPr = r._r.get_or_add_rPr()
    # w:spacing is character tracking in twentieths of a point
    rPr.append(rPr.makeelement(qn("w:spacing"),
                               {qn("w:val"): str(round(tracking_em * 7.5 * 20))}))


def build_docx(out, *, letterhead, korean):
    doc = docx.Document()
    sec = doc.sections[0]
    sec.top_margin = sec.bottom_margin = Mm(18)
    sec.left_margin = sec.right_margin = Mm(18)

    style = doc.styles["Normal"]
    style.font.name = LATIN
    style.font.size = DocPt(10.5)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), HANGUL)

    if letterhead:
        name = "이한보람" if korean else NAME
        affil = ("울산과학기술원(UNIST) 반도체소재부품대학원 교수" if korean
                 else "Professor, Graduate School of Semiconductor Materials & Devices "
                      "Engineering,\nUlsan National Institute of Science and Technology (UNIST)")
        table = doc.add_table(rows=1, cols=2)
        table.autofit = False
        left, right = table.rows[0].cells
        left.width, right.width = Cm(2.6), Cm(14.8)
        left.paragraphs[0].add_run().add_picture(
            str(ASSETS / "unist-emblem.png"), width=Cm(2.3))
        p = right.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.paragraph_format.space_after = DocPt(3)
        dset(p.add_run(name), 24, D_NAVY)
        for line, size, colour in (
            (affil, 9.5, D_NAVY),
            ("052-217-3218, hbrlee@unist.ac.kr, https://hbrl-research.group", 8.5, D_MUTED),
            ("Executive Editor, Chemistry of Materials, ACS Publications", 9.5, D_NAVY),
            ("(919) 650-1459, (202) 350-3293, lee-office@cm.acs.org", 8.5, D_MUTED),
        ):
            q = right.add_paragraph()
            q.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            q.paragraph_format.space_after = DocPt(1)
            dset(q.add_run(line), size, colour)
    else:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = DocPt(2)
        p.add_run().add_picture(str(ASSETS / "logo/hbrlrg-horizontal.png"), width=Cm(3.4))
        q = doc.add_paragraph()
        q.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        dset(q.add_run("https://hbrl-research.group · hbrlee@unist.ac.kr"), 8.5, D_MUTED)

    rule(doc)

    if letterhead:
        body = [("2026. 00. 00.", 10.5, D_MUTED, False),
                ("받는 분 / Recipient" if korean else "Recipient", 10.5, D_INK, True),
                ("안녕하십니까." if korean else "Dear …,", 10.5, D_INK, False),
                ("본문을 여기에 씁니다. 단락 사이는 10pt, 들여쓰기는 하지 않습니다."
                 if korean else
                 "Body text goes here. Ten points between paragraphs, no indent.",
                 10.5, D_INK, False),
                ("감사합니다." if korean else "Sincerely,", 10.5, D_INK, False),
                ("", 10.5, D_INK, False),
                ("", 10.5, D_INK, False),
                (NAME, 10.5, D_NAVY, True)]
    else:
        body = [("문서 제목" if korean else "Document title", 20, D_NAVY, True),
                ("작성자 · 2026. 00. 00. · v1" if korean else
                 "Author · 00 Month 2026 · v1", 8.5, D_MUTED, False),
                ("1. 첫 번째 섹션" if korean else "1. First section", 13, D_NAVY, True),
                ("본문을 여기에 씁니다. 섹션마다 번호를 붙이고, 표는 3–5열 이내로 유지합니다."
                 if korean else
                 "Body text goes here. Number every section; keep tables to three "
                 "to five columns.", 10.5, D_INK, False)]

    for txt, size, colour, bold in body:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = DocPt(10)
        if txt:
            dset(p.add_run(txt), size, colour, bold=bold)

    if letterhead:
        doc_footer(doc, "052-217-3218  ·  hbrlee@unist.ac.kr  ·  https://hbrl-research.group",
                   0.05)
    else:
        doc_footer(doc, "HBRL RESEARCH GROUP  ·  GRADUATE SCHOOL OF SEMICONDUCTOR "
                        "MATERIALS & DEVICES ENGINEERING, UNIST", 0.07)

    doc.save(out)


def main():
    outdir = Path(sys.argv[1] if len(sys.argv) > 1 else "office")
    outdir.mkdir(parents=True, exist_ok=True)

    n = build_pptx(outdir / "hbrlrg-slides.pptx")
    print(f"  hbrlrg-slides.pptx        {n} slides")

    for name, lh, ko in (("letterhead-kr", True, True), ("letterhead-en", True, False),
                         ("plain-kr", False, True), ("plain-en", False, False)):
        build_docx(outdir / f"hbrlrg-{name}.docx", letterhead=lh, korean=ko)
        print(f"  hbrlrg-{name}.docx")


if __name__ == "__main__":
    main()
