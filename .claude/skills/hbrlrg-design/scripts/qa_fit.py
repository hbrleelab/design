"""Geometric QA for a .pptx without a renderer.

When LibreOffice is unavailable, text fit can still be checked by measuring
against the real Archivo advance widths rather than eyeballing a thumbnail.
Reports text taller than its box, shapes off-slide, and text-box collisions.

    python3 qa_fit.py deck.pptx [dir-with-Archivo-Regular.ttf-and-Bold.ttf]

Run it against the original file too: that separates defects you introduced
from ones the source already had.
"""
import os
import sys

from fontTools.ttLib import TTFont
from pptx import Presentation
from pptx.util import Emu, Pt

FONT_DIR = os.environ.get("HBRLRG_FONT_DIR", "fonts")  # dir holding Archivo-*.ttf
SLIDE_W, SLIDE_H = 13.333, 7.5
DEFAULT_LINE = 1.2  # PowerPoint's single line spacing


class Metrics:
    def __init__(self, path):
        self.font = TTFont(path, lazy=True)
        self.upm = self.font["head"].unitsPerEm
        self.cmap = self.font.getBestCmap()
        self.hmtx = self.font["hmtx"]
        self.order = self.font.getGlyphOrder()
        self._cache = {}

    def advance(self, ch):
        if ch in self._cache:
            return self._cache[ch]
        name = self.cmap.get(ord(ch))
        if name is None:
            name = self.cmap.get(ord("?"), self.order[0])
        adv = self.hmtx[name][0] / self.upm
        self._cache[ch] = adv
        return adv

    def width(self, text, size_pt):
        return sum(self.advance(c) for c in text) * size_pt


def wrap_lines(text, size_pt, avail_pt, metrics):
    """Greedy word wrap; returns the line count."""
    if not text.strip():
        return 1
    lines = 0
    for hard in text.split("\n"):
        words = hard.split(" ")
        cur = ""
        n = 1
        for w in words:
            trial = w if not cur else cur + " " + w
            if metrics.width(trial, size_pt) <= avail_pt or not cur:
                cur = trial
            else:
                n += 1
                cur = w
        lines += n
    return lines


def main(path, font_dir=None):
    font_dir = font_dir or FONT_DIR
    reg = Metrics(f"{font_dir}/Archivo-Regular.ttf")
    bold = Metrics(f"{font_dir}/Archivo-Bold.ttf")

    prs = Presentation(path)
    overflow, offslide, collide = [], [], []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            L, T = Emu(sh.left).inches, Emu(sh.top).inches
            W, H = Emu(sh.width).inches, Emu(sh.height).inches
            if L < -0.01 or T < -0.01 or L + W > SLIDE_W + 0.01 or T + H > SLIDE_H + 0.01:
                offslide.append((idx, sh.shape_id, round(L, 2), round(T, 2),
                                 round(L + W, 2), round(T + H, 2)))
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue

            tf = sh.text_frame
            ml = Emu(tf.margin_left).inches if tf.margin_left is not None else 0.1
            mr = Emu(tf.margin_right).inches if tf.margin_right is not None else 0.1
            mt = Emu(tf.margin_top).inches if tf.margin_top is not None else 0.05
            mb = Emu(tf.margin_bottom).inches if tf.margin_bottom is not None else 0.05
            avail_pt = max((W - ml - mr) * 72, 1)

            total_pt = 0.0
            for para in tf.paragraphs:
                if not para.runs:
                    total_pt += 12 * DEFAULT_LINE
                    continue
                size = max(
                    (r.font.size.pt for r in para.runs if r.font.size), default=18.0
                )
                is_bold = any(r.font.bold for r in para.runs)
                m = bold if is_bold else reg
                text = "".join(r.text for r in para.runs)
                if tf.word_wrap is False:
                    n = 1
                else:
                    n = wrap_lines(text, size, avail_pt, m)
                spacing = para.line_spacing if isinstance(para.line_spacing, float) else DEFAULT_LINE
                total_pt += n * size * spacing
                if para.space_after is not None:
                    total_pt += para.space_after.pt
                if para.space_before is not None:
                    total_pt += para.space_before.pt

            need = total_pt / 72 + mt + mb
            if need > H + 0.02:
                overflow.append((idx, round(need, 2), round(H, 2), round(need - H, 2),
                                 tf.text.strip().replace("\n", " / ")[:52]))
            boxes.append((L, T, W, H, tf.text.strip()[:26]))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
                oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                if ox > 0.05 and oy > 0.05:
                    collide.append((idx, a[4], b[4], round(ox, 2), round(oy, 2)))

    print(f"=== 텍스트가 상자보다 큼 ({len(overflow)}건) ===")
    for s, need, have, d, t in overflow:
        print(f"  S{s:2}  필요 {need:.2f}in > 상자 {have:.2f}in  (+{d:.2f})  {t!r}")
    print(f"\n=== 슬라이드 밖으로 나감 ({len(offslide)}건) ===")
    for s, sid, l, t, r, b in offslide:
        print(f"  S{s:2}  id={sid} x {l}~{r}  y {t}~{b}")
    print(f"\n=== 텍스트 상자 겹침 ({len(collide)}건) ===")
    for s, a, b, ox, oy in collide:
        print(f"  S{s:2}  {a!r} ∩ {b!r}  {ox}x{oy}in")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit("usage: qa_fit.py deck.pptx [font-dir-with-Archivo-ttf]")
    if len(sys.argv) > 2:
        main(sys.argv[1], sys.argv[2])
    else:
        main(sys.argv[1])
